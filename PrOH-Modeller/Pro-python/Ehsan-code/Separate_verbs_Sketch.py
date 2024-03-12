import csv
import nltk
from nltk.tokenize import word_tokenize
from nltk import pos_tag
from pptx import Presentation
from pptx.util import Inches
import pydot

# Download NLTK data if not already installed
nltk.download('punkt')
nltk.download('averaged_perceptron_tagger')

# Input CSV file name
input_file = 'sipoc_table.csv'  # Replace with the path to your input CSV file
output_file = 'verbs.csv'
pptx_file = 'verbs.pptx'

# List to store verbs
verbs_data = []

def is_verb(tag):
    return tag in ['VB', 'VBD', 'VBG', 'VBN', 'VBP', 'VBZ']

# Create a PPTX presentation
prs = Presentation()

# Create a single slide for verb visualization
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

with open(input_file, 'r', newline='') as csv_input:
    reader = csv.reader(csv_input)

    for row in reader:
        for cell in row:
            # Split the cell content into individual words
            words = word_tokenize(cell)
            
            # Part-of-speech tagging
            tagged_words = pos_tag(words)
            
            # Extract verbs
            verbs = [word for word, tag in tagged_words if is_verb(tag)]
            verbs_data.extend(verbs)

# Write verbs to the output CSV file
with open(output_file, 'w', newline='') as csv_output:
    writer = csv.writer(csv_output)
    writer.writerow(['Verb'])  # Header
    writer.writerows([[verb] for verb in verbs_data])

print(f"Verbs have been saved to '{output_file}'")

# Create a graphical representation of verbs using pydot
#graph = pydot.Dot(graph_type='graph')
#for verb in verbs_data:
 #   node = pydot.Node(verb)
  #  graph.add_node(node)

# Save the PPTX file with verbs as nodes
#image_path = 'verbs_graph.png'
#graph.write_png(image_path)

# Calculate the positions to fit all nodes within the slide
num_nodes = len(verbs_data)
max_nodes_per_row = 5
num_rows = (num_nodes + max_nodes_per_row - 1) // max_nodes_per_row
node_width = Inches(1.5)
node_height = Inches(0.5)
left_margin = Inches(0.5)
top_margin = Inches(1)

# Create shapes in PowerPoint
left = left_margin
top = top_margin
for verb in verbs_data:
    if left + node_width > Inches(10):
        left = left_margin
        top += node_height
    slide.shapes.add_textbox(left, top, node_width, node_height).text_frame.text = verb
    left += node_width

# Save the PPTX file
prs.save(pptx_file)
print(f"PPTX file with verbs has been created: '{pptx_file}'")