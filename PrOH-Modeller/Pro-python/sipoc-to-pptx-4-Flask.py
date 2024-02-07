from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import csv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import spacy




# Load the spaCy English model
nlp = spacy.load("en_core_web_sm")

# Read the CSV file and extract the second row
csv_filename = 'uploads/sipoc_table.csv'  # Replace with the path to your CSV file
with open(csv_filename, 'r') as csv_file:
    csv_reader = csv.reader(csv_file)
    rows = list(csv_reader)

# Create a PowerPoint presentation
presentation = Presentation()

# Create a blank slide
slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Define position and size for oval shapes
num_cells = len(rows[1])
oval_width = Inches(2)
oval_height = Inches(0.8)
text_font_size = Pt(18)
slide_width = presentation.slide_width
slide_height = presentation.slide_height

# Calculate the diagonal line positions
diagonal_length = (slide_width ** 2 + slide_height ** 2) ** 0.5
x_step = slide_width / num_cells
y_step = slide_height / num_cells

# Create oval shapes and add text to them with specified font and background colors
for i, cell in enumerate(rows[1]):
    left = i * x_step
    top = i * y_step
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)

    if i % 2 == 0:  # Even columns (0-based index)
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(255, 0, 0)	  # Red background
    else:
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background

    text_frame = oval.text_frame
    text_frame.text = cell
    text_frame.paragraphs[0].font.size = text_font_size
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Use spaCy to extract verbs from the second row
verbs = []
for cell in rows[1]:
    doc = nlp(cell)
    verbs.extend([token.text for token in doc if token.pos_ == "VERB"])

# Create nodes for verbs with no border and background at the bottom of the slide
for i, verb in enumerate(verbs):
    left = i * x_step
    top = slide_height - oval_height - Inches(0.2)
    verb_oval = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
    verb_oval.line.fill.solid()
    verb_oval.line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # No border
    verb_oval.line.width = Pt(0)  # No border width
    verb_oval.shadow.inherit = False  # No shadow
    verb_oval.fill.solid()
    verb_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # No background
    verb_text_frame = verb_oval.text_frame
    p = verb_text_frame.add_paragraph()
    p.text = verb
    p.font.size = text_font_size
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
    p.alignment = PP_ALIGN.CENTER

# Save the PowerPoint presentation
pptx_filename = 'uploads/output_presentation4.pptx'
presentation.save(pptx_filename)

print(f'Second row of CSV file has been converted to an editable PowerPoint presentation: "{pptx_filename}"')



