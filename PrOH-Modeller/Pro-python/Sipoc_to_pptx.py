from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import csv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import spacy
import logging
import os
import sys

# Create a new logger instance for Seperate_verbs.py
logger = logging.getLogger('sipoc_to_pptx_logger')
logger.setLevel(logging.INFO)  # Set the log level to INFO

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('sipoc_to_pptx.log')
file_handler.setLevel(logging.INFO)  # Set the log level for the handler to INFO

# Define a formatter for the log messages
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)  # Attach the formatter to the handler

# Add the file handler to the logger
logger.addHandler(file_handler)

try:
    # Get the file name with the identifier from the command-line arguments
    filename_with_identifier = sys.argv[1]
    logger.info("Looking for the filepath")

    filename_without_extension = os.path.splitext(filename_with_identifier)[0]
    logger.info(f"Filename without extension: {filename_without_extension}")

    # Construct the complete file path
    file_path = os.path.join(filename_with_identifier)
    logger.info(f"Uploaded filepath: {file_path}")

    with open(file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        rows = list(csv_reader)

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Create a blank slide
    slide_layout = presentation.slide_layouts[5]  # Blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Define position and size for oval shapes
    num_cells = 7
    oval_width = Inches(2)
    oval_height = Inches(0.8)
    text_font_size = Pt(18)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Calculate the diagonal line positions
    #diagonal_length = (slide_width ** 2 + slide_height ** 2) ** 0.5
    x_step = slide_width / num_cells
    y_step = slide_height / num_cells

    # Create oval shapes with red background for the first cells in columns A and F
    for i, cell in enumerate(rows[1]):
        # Skip column D
        if i == 3:
            continue
        left = i * x_step
        top = i * y_step
        
        if i == 0:  #First cell in column A
            # Extract the first word from cells A1 and F1
            red_text = rows[1][0].split()[0]  # Get the first word from cell A1
            # Create rectangle shape
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for rectangles

            # Create oval shape with red background
            oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left+x_step, top+y_step, oval_width, oval_height)
            oval_shape.fill.solid()
            oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
    
        elif i == 5:  # First cell in column F
            red_text = rows[1][5].split()[0]  # Get the first word from cell F1
            
            # Create rectangle shape
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left+x_step, top+y_step, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for rectangles
            
            # Create oval shape with red background
            oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
            oval_shape.fill.solid()
            oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
            
        elif i == 1: # Column B
            # Create oval shape with green background
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left+x_step, top+y_step, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background
        elif i == 2:  # Column C
            # Create oval shape with red background
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left+x_step, top+y_step, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
        elif i == 4:  # Column E
            # Create oval shape with green background
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background

        text_frame = shape.text_frame
        text_frame.text = cell
        text_frame.paragraphs[0].font.size = text_font_size
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Add text to oval shape
        oval_text_frame = oval_shape.text_frame
        oval_text_frame.text = red_text  # Use the extracted first word as text
        oval_text_frame.paragraphs[0].font.size = text_font_size
        oval_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
        oval_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    # Load the spaCy English model
    nlp = spacy.load("en_core_web_sm")

    # Use to extract verbs from the second row, ignoring column D
    verbs = []
    for i, cell in enumerate(rows[1]):
        # Skip column D
        if i == 3:
            continue
    
        # Apply spaCy NLP pipeline to the cell content
        doc = nlp(cell)
        # Extract verbs
        for token in doc:
            if token.pos_ == 'VERB':
                verbs.append(token.text)

    # Divide the list of verbs into chunks of five
    verbs_chunks = [verbs[i:i+5] for i in range(0, len(verbs), 5)]

    # Define position and size for text boxes
    left_margin = Inches(0.5)
    top_margin = presentation.slide_height - Inches(0.5)  # Place at the bottom of the slide
    box_width = Inches(1.0)  # Width of each text box
    box_height = Inches(0.5)  # Height of each text box

    # Create text boxes for each chunk of verbs and add them to the bottom of the slide
    for chunk in verbs_chunks:
        left = left_margin
        for word in chunk:
            textbox = slide.shapes.add_textbox(left, top_margin, box_width, box_height)
            textbox.text_frame.text = word
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Align the text to the center
            left += box_width
        top_margin -= box_height  # Move to the next line

    # Save the PowerPoint presentation with the same identifier
    pptx_filename = os.path.join(filename_without_extension + '.pptx')
    presentation.save(pptx_filename)
    logger.info("Output presentation is saved successfully.")
    logger.info(f"Output presentation path: {os.path.join(filename_without_extension + '.pptx')}")
    print(f'Second row of CSV file has been converted to an editable PowerPoint presentation: "{pptx_filename}"')

except Exception as e:
    # Log any exceptions that occur
    logging.error(f'An error occurred: {str(e)}')
