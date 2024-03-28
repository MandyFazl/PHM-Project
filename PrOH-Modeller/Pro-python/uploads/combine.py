from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def combine_pptx_files(input_files, output_file):
    # Create a presentation object for the combined output
    combined_presentation = Presentation()

    for file_path in input_files:
        # Open each input presentation
        input_presentation = Presentation(file_path)

        # Iterate through each slide in the input presentation
        for slide in input_presentation.slides:
            # Create a new slide in the combined presentation and copy the content from the input slide
            new_slide = combined_presentation.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                # Check if the shape is an auto shape
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Add the shape to the new slide
                    new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                    # Copy the text
                    new_shape.text = shape.text
                    # Copy the formatting properties if they are defined
                    try:
                        if shape.fill.fore_color.rgb is not None:
                            new_shape.fill.solid()
                            new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                        if shape.line.color.rgb is not None:
                            new_shape.line.color.rgb = shape.line.color.rgb
                        if shape.line.width is not None:
                            new_shape.line.width = shape.line.width
                    except AttributeError:
                        pass
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    # Add a text box to the new slide
                    new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    # Copy text and formatting
                    for paragraph in shape.text_frame.paragraphs:
                        new_paragraph = new_shape.text_frame.add_paragraph()
                        new_paragraph.text = paragraph.text
                        # Set font color to black for each run
                        for run in new_paragraph.runs:
                            run.font.color.rgb = RGBColor(0, 0, 0)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Handle grouped shapes
                    for grouped_shape in shape.shapes:
                        if grouped_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                            # Add a text box to the new slide
                            new_shape = new_slide.shapes.add_textbox(grouped_shape.left, grouped_shape.top, grouped_shape.width, grouped_shape.height)
                            # Copy text and formatting
                            for paragraph in grouped_shape.text_frame.paragraphs:
                                new_paragraph = new_shape.text_frame.add_paragraph()
                                new_paragraph.text = paragraph.text
                                # Set font color to black for each run
                                for run in new_paragraph.runs:
                                    run.font.color.rgb = RGBColor(0, 0, 0)

    # Save the combined presentation to the output file
    combined_presentation.save(output_file)

if __name__ == "__main__":
    # List of input PPTX files
    input_files = [
    'make_car_d27a592c_2407a060.pptx',
    'make_car_d27a592c_2407a060_non-cp-statement.pptx',
    'make_car_d27a592c_2407a060_verbs.pptx',
    'make_car_d27a592c_2407a060_subbubbles.pptx',
    'make_car_d27a592c_2407a060_decision-bubbles.pptx'
    ]
    # Output file name for the combined presentation
    output_file = 'combined_output.pptx'

    # Combine the input PPTX files into a single output file
    combine_pptx_files(input_files, output_file)

    print(f"Combined presentation saved to {output_file}")
