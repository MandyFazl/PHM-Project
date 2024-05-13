import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Input CSV file name
input_file = 'parentheses_words.csv'  # Replace with the path to your input CSV file
output_pptx_file = 'parentheses_words-unique.pptx'

# Set to store unique words between parentheses
unique_parentheses_words = set()

# Read unique words between parentheses from the CSV file
with open(input_file, 'r', newline='') as csv_input:
    reader = csv.reader(csv_input)
    next(reader)  # Skip the header row

    for row in reader:
        for cell in row:
            unique_parentheses_words.add(cell)

# Create a PowerPoint presentation
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout

# Create a single slide with a white background
slide = prs.slides.add_slide(slide_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Calculate the positions to fit all nodes within the slide
num_nodes = len(unique_parentheses_words)
node_width = Inches(2.5)
node_height = Inches(1.0)
left_margin = Inches(0.5)
top_margin = Inches(1.0)
left = left_margin
top = top_margin

# Create elliptical shapes in PowerPoint for unique words
for word in unique_parentheses_words:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    text = shape.text_frame.add_paragraph()
    text.text = word
    text.alignment = 1  # Center align text
    text.font.size = Pt(18)
    text.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
    shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    shape.line.width = Pt(1.5)  # Border width
    left += node_width

    if left + node_width > Inches(10):
        left = left_margin
        top += node_height

# Save the PowerPoint file
prs.save(output_pptx_file)
print(f"PPTX file with editable elliptical nodes for unique values has been created: '{output_pptx_file}'")
