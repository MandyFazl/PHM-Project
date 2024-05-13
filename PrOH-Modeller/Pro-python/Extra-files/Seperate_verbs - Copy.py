import csv
from nltk.tokenize import word_tokenize
from nltk import pos_tag
from pptx import Presentation
from pptx.util import Inches
import os
import sys
import logging
from pptx.enum.text import PP_ALIGN

# Configure logging
logging.basicConfig(filename='separate_verbs.log', level=logging.INFO)

filename_with_identifier = sys.argv[1]
filename_without_extension = os.path.splitext(filename_with_identifier)[0]
file_path = os.path.join(filename_with_identifier)

temp_csv_filename = filename_without_extension + '_verbs.csv'

try:
    with open(file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        rows = list(csv_reader)
        # Remove first two rows
        rows = rows[2:]
        # Read the CSV file and extract verbs and their respective cells
        
        
        verb_data = []  # Use spaCy to extract verbs
        j = 0
        while j< len(rows):
            print(f'length: "{len(rows)}"')
            row=rows[j]  
            for i, cell in enumerate(row):
                if  i==0 or i==5:
                    continue  
    # with open(file_path, 'r') as csv_file:
    #     csv_reader = csv.reader(csv_file)
    #     for row in csv_reader:
    #         for idx, cell in enumerate(row):
    #             # Skip columns A and F (index 0 and 5)
    #             if idx in [0, 5]:
    #                 continue

                # Split the cell content into individual words
                words = word_tokenize(cell)
                # Part-of-speech tagging
                tagged_words = pos_tag(words)
                # Check if the cell contains a verb
                if any(tag.startswith('VB') for _, tag in tagged_words):
                    verb_data.append((cell,))
            
            # Write verbs and their respective cells to the temporary CSV file
            with open(temp_csv_filename, 'w', newline='') as temp_csv_file:
                writer = csv.writer(temp_csv_file)
                writer.writerow(['Cell'])  # Header
                writer.writerows(verb_data)
            j += 1

    # Create a PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Calculate the positions to fit all cells within the slide
    max_cells_per_row = 5
    num_cells = len(verb_data)
    num_rows = (num_cells + max_cells_per_row - 1) // max_cells_per_row
    cell_width = Inches(1.5)
    cell_height = Inches(0.5)
    left_margin = Inches(0.5)
    top_margin = Inches(1)

    # Create shapes in PowerPoint
    left = left_margin
    top = top_margin
    for cell_content in verb_data:
        cell = cell_content[0]
        if left + cell_width > Inches(10):
            left = left_margin
            top += cell_height
        textbox = slide.shapes.add_textbox(left, top, cell_width, cell_height)
        textbox.text_frame.text = cell
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
        left += cell_width

    # Save the PowerPoint presentation
    output_pptx_file = filename_without_extension + '_verbs.pptx'
    prs.save(output_pptx_file)
    print(f"PPTX file with cells containing verbs has been created: '{output_pptx_file}'")

except Exception as e:
    # Log any exceptions that occur
    logging.error(f'An error occurred: {str(e)}')
