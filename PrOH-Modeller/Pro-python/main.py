from flask import Flask, render_template, request, send_file, jsonify, session, make_response
import subprocess
import os
import logging
import uuid  # for generating unique identifiers
import zipfile
from pptx import Presentation



app = Flask(__name__)
app.secret_key = os.urandom(24)  # Secret key for session

# Configure logging
logging.basicConfig(filename='app.log', level=logging.INFO)

# Define upload status constants
UPLOAD_SUCCESS = "File uploaded and processed successfully"
NO_FILE_SELECTED = "No file selected"


@app.route("/")
def index():
    return render_template('index.html', upload_status="")


@app.route('/upload', methods=['POST'])
def upload():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file selected'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400

        # Generate a unique identifier
        unique_identifier = str(uuid.uuid4())[:8]  # Adjust the length of the unique code as needed

        # Get the filename and extension
        filename, extension = os.path.splitext(file.filename)

        # Append the unique identifier to the filename
        filename_with_identifier = f"{filename}_{unique_identifier}{extension}"

        # Save the file with the new filename
        file_path = os.path.join('uploads', filename_with_identifier)
        file.save(file_path)
        logging.info("CSV file saved successfully.")

        # Store the file path in session
        session['file_path'] = file_path

        logging.info(f"Uploaded filename: {file.filename}")
        logging.info(f"Uploaded filepath: {file_path}")

        # Construct the filename without the extension
        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        logging.info(f"Filename without extension: {filename_without_extension}")

        # Call your python scripts with the uploaded file as an argument
        logging.info("calling your python scripts.")
        subprocess.run(['python', 'Sipoc_to_pptx.py', file_path])
        logging.info("Sipoc_to_pptx.py executed successfully.")
        subprocess.run(['python', 'Non_Core_Process_Statement.py', file_path])
        logging.info("Non_Core_Process_Statement.py executed successfully.")
        subprocess.run(['python', 'Seperate_verbs.py', file_path])
        logging.info("Seperate_verbs.py executed successfully.")
        subprocess.run(['python', 'Bracket.py', file_path])
        logging.info("Bracket.py executed successfully.")
        subprocess.run(['python', 'Decision_Bubbles.py', file_path])
        logging.info("Decision_Bubbles.py executed successfully.")

        return jsonify({'message': 'Upload successful'}), 200
    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred: {str(e)}')
        return jsonify({'error': 'An error occurred while processing the file'}), 500



@app.route('/download_all_files')
def download_all_files():
    try:
        file_path = session.get('file_path')
        if not file_path:
            return "No file to download"

        filename_with_identifier = os.path.basename(file_path)
        filename_without_extension = os.path.splitext(filename_with_identifier)[0]

        pptx_path = os.path.join('uploads', f'{filename_without_extension}.pptx')
        non_cp_statement_path = os.path.join('uploads', f'{filename_without_extension}_non-cp-statement.pptx')
        verbs_path = os.path.join('uploads', f'{filename_without_extension}_verbs.pptx')
        decision_bubbles_path = os.path.join('uploads', f'{filename_without_extension}_decision-bubbles.pptx')
        subbubbles_path = os.path.join('uploads', f'{filename_without_extension}_subbubbles.pptx')

        logging.info(f"pptx_path: {pptx_path}")
        logging.info(f"non_cp_statement_path: {non_cp_statement_path}")
        logging.info(f"verbs_path: {verbs_path}")
        logging.info(f"decision_bubbles_path: {decision_bubbles_path}")
        logging.info(f"subbubbles_path: {subbubbles_path}")

        combined_presentation = Presentation()

        for file_path in [pptx_path, non_cp_statement_path, verbs_path, decision_bubbles_path, subbubbles_path]:
            if os.path.exists(file_path):
                logging.info(f"Adding slide from: {file_path}")
                slide_layout = combined_presentation.slide_layouts[5]  # Assuming a blank slide layout
                slide = combined_presentation.slides.add_slide(slide_layout)
                slide.shapes.add_picture(file_path, left=0, top=0, width=None, height=None)

        combined_presentation_path = os.path.join('uploads', f'{filename_without_extension}_combined.pptx')
        combined_presentation.save(combined_presentation_path)
        logging.info(f"combined_presentation_path: {combined_presentation_path}")
        logging.info("combined_presentation_path saved")

        response = make_response(send_file(combined_presentation_path, as_attachment=True))
        response.headers['Content-Disposition'] = 'attachment; filename=combined_presentation.pptx'
        return response

    except Exception as e:
        # Handle exceptions
        logging.error(f'An error occurred while downloading all files: {str(e)}')
        return "An error occurred while downloading all files."


@app.route('/download_pptx')
def download_pptx():
    try:
        # Get the file path from session
        file_path = session.get('file_path')
        if file_path:
            # Specify the path to the PowerPoint file
            filename_with_identifier = os.path.basename(file_path)
            filename_without_extension = os.path.splitext(filename_with_identifier)[0]
            pptx_filename = os.path.join('uploads', f'{filename_without_extension}.pptx')
            logging.info("PowerPoint file downloaded successfully.")

            # Return the file as an attachment
            return send_file(pptx_filename, as_attachment=True)
        else:
            return "No file to download"
    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred: {str(e)}')
        # Return an error message or redirect to an error page
        return "An error occurred while downloading the file."
    
@app.route('/download_NC_Statement')
def download_NC_Statement():
    try:
        # Get the file path from session
        file_path = session.get('file_path')
        if file_path:
            # Specify the path to the PowerPoint file
            filename_with_identifier = os.path.basename(file_path)
            filename_without_extension = os.path.splitext(filename_with_identifier)[0]
            Non_cp_filename = os.path.join('uploads', filename_without_extension +'_non-cp-statement'+'.pptx')
            logging.info("Non-CP-Statement file downloaded successfully.")

            # Return the file as an attachment
            return send_file(Non_cp_filename, as_attachment=True)
        else:
            return "No file to download"
    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred: {str(e)}')
        # Return an error message or redirect to an error page
        return "An error occurred while downloading the file."


@app.route('/download_verbs')
def download_verbs():
    try:
        # Specify the path to the verbs.csv file
        file_path = session.get('file_path')
        if file_path:
            filename_with_identifier = os.path.basename(file_path)
            filename_without_extension = os.path.splitext(filename_with_identifier)[0]
            verbs_pptxfile= os.path.join('uploads',filename_without_extension +'_verbs'+'.pptx')
            logging.info("Verbs pptx file downloaded successfully.")


        # Return the file as an attachment
        return send_file(verbs_pptxfile, as_attachment=True)
    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred: {str(e)}')
        # Return an error message or redirect to an error page
        return "An error occurred while downloading the file."
    
@app.route('/download_decision_bubbles')
def download_decision_bubbles():
        try:
            # Get the file path from session
            file_path = session.get('file_path')
            if file_path:
                # Specify the path to the PowerPoint file
                filename_with_identifier = os.path.basename(file_path)
                filename_without_extension = os.path.splitext(filename_with_identifier)[0]
                pptx_filename = os.path.join('uploads', filename_without_extension + '_decision-bubbles' + '.pptx')
                logging.info("decision_bubbles pptx file downloaded successfully.")

                # Return the file as an attachment
                return send_file(pptx_filename, as_attachment=True)
            else:
                return "No file to download"
        except Exception as e:
            # Log any exceptions that occur
            logging.error(f'An error occurred: {str(e)}')
            # Return an error message or redirect to an error page
            return "An error occurred while downloading the file."
    
    
@app.route('/download_Sub_bubbles')
def download_Sub_bubbles():
    try:
        # Specify the path to the verbs.csv file
        file_path = session.get('file_path')
        if file_path:
            filename_with_identifier = os.path.basename(file_path)
            filename_without_extension = os.path.splitext(filename_with_identifier)[0]
            verbs_pptxfile= os.path.join('uploads',filename_without_extension +'_subbubbles'+'.pptx')
            logging.info(f"Sub_bubbles pptx file path {verbs_pptxfile}")
            logging.info("Sub_bubbles pptx file downloaded successfully.")


        # Return the file as an attachment
        return send_file(verbs_pptxfile, as_attachment=True)
    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred: {str(e)}')
        # Return an error message or redirect to an error page
        return "An error occurred while downloading the file."
    

    

if __name__ == '__main__':

    app.run(host="0.0.0.0")
