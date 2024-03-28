import csv
import logging
import os
import sys
from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import spacy
from merge2functions import process_script1, process_script2

# Create a logger instance
logger = logging.getLogger('merged_logger')
logger.setLevel(logging.INFO)

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('merged_script.log')
file_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)
logger.addHandler(file_handler)

# Load the spaCy English model
nlp = spacy.load("en_core_web_sm")

# Create a single Presentation object
combined_presentation = Presentation()

def process_script1(filename_with_identifier):
    try:
        logger.info("Processing Script 1")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        logger.info(f"Filename without extension: {filename_without_extension}")

        # Construct the complete file path
        file_path = os.path.join(filename_with_identifier)
        logger.info(f"Uploaded filepath: {file_path}")

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)

        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Define position and size for oval shapes
        num_cells = len(rows[1])
        oval_width = Inches(2)
        oval_height = Inches(0.8)
        text_font_size = Pt(18)
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

        # Calculate the diagonal line positions
        diagonal_length = (slide_width ** 2 + slide_height ** 2) ** 0.5
        x_step = slide_width / num_cells
        y_step = slide_height / num_cells

        # Create oval shapes and add text to them with specified font and background colors
        for i, cell in enumerate(rows[1]):
            left = i * x_step
            top = i * y_step
            if i == 0 or i == 5:  # First cell in column A and F
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for rectangles
            elif i == 1 or i == 4:  # Column B and E
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background
            elif i == 2:  # Column C
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background

            text_frame = shape.text_frame
            text_frame.text = cell
            text_frame.paragraphs[0].font.size = text_font_size
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Use spaCy to extract verbs from the second row
        verbs = []
        for cell in rows[1]:
            doc = nlp(cell)
            verbs.extend([token.text for token in doc if token.pos_ == "VERB"])

        # Create nodes for verbs with no border and background at the bottom of the slide
        for i, verb in enumerate(verbs):
            left = i * x_step
            top = slide_height - oval_height - Inches(0.2)
            verb_oval = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
            verb_oval.line.fill.solid()
            verb_oval.line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # No border
            verb_oval.line.width = Pt(0)  # No border width
            verb_oval.shadow.inherit = False  # No shadow
            verb_oval.fill.solid()
            verb_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # No background
            verb_text_frame = verb_oval.text_frame
            p = verb_text_frame.add_paragraph()
            p.text = verb
            p.font.size = text_font_size
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            p.alignment = PP_ALIGN.CENTER

        logger.info("Script 1 executed successfully.")

    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred in Script 1: {str(e)}')

def process_script2(filename_with_identifier):
    try:
        logger.info("Processing Script 2")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        logger.info(f"Filename without extension: {filename_without_extension}")

        # Construct the complete file path
        file_path = os.path.join(filename_with_identifier)
        logger.info(f"Uploaded filepath: {file_path}")

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)

        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]  # Blank slide layout
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Define initial positions for ovals
        left = Inches(0.5)
        top = Inches(5)

        # Iterate over all cells in the CSV file
        for row in rows:
            for idx, cell in enumerate(row):
                # Skip processing cells in columns A and F (indexes 0 and 5 respectively)
                if idx in [0, 5]:
                    continue

                # Check if the cell begins with "or:" case-insensitive
                if cell.lower().startswith("or:"):
                    # Extract the phrase after "or:"
                    phrase = cell[3:].strip()

                    # Check if the phrase contains parentheses
                    if '(' in phrase and ')' in phrase:
                        # Split the phrase into main and inner phrases
                        main_phrase, inner_phrase = phrase.split('(', 1)
                        inner_phrase = inner_phrase[:-1]  # Remove the closing parenthesis

                        # Create the outer oval for the main phrase
                        outer_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(3), Inches(1.5))
                        outer_oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                        outer_oval.line.width = Pt(1)  # Border width
                        outer_oval.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dotted line border
                        outer_oval.fill.solid()
                        outer_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                        # Add the main phrase to the outer oval
                        text_frame_outer = outer_oval.text_frame
                        p_outer = text_frame_outer.add_paragraph()
                        p_outer.text = main_phrase.strip()
                        p_outer.font.size = Pt(14)  # Adjust font size as needed
                        p_outer.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                        p_outer.alignment = PP_ALIGN.CENTER  # Center align text

                        # Create the inner oval for the phrase inside parentheses
                        inner_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.5), top + Inches(0.3), Inches(2), Inches(0.8))
                        inner_oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                        inner_oval.line.width = Pt(1)  # Border width
                        inner_oval.fill.solid()
                        inner_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                        # Add the inner phrase to the inner oval
                        text_frame_inner = inner_oval.text_frame
                        p_inner = text_frame_inner.add_paragraph()
                        p_inner.text = inner_phrase.strip()
                        p_inner.font.size = Pt(12)  # Adjust font size as needed
                        p_inner.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                        p_inner.alignment = PP_ALIGN.CENTER  # Center align text

                    else:
                        # Create a black dotted line oval for the extracted phrase
                        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(3), Inches(1.5))
                        oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                        oval.line.width = Pt(1)  # Border width
                        oval.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dotted line border
                        oval.fill.solid()
                        oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                        # Add the extracted phrase to the oval
                        text_frame = oval.text_frame
                        p = text_frame.add_paragraph()
                        p.text = phrase.strip()
                        p.font.size = Pt(14)  # Adjust font size as needed
                        p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                        p.alignment = PP_ALIGN.CENTER  # Center align text

                    # Update positions for the next oval
                    left += Inches(3.5)
                    if left + Inches(3) > Inches(10):
                        left = Inches(0.5)
                        top += Inches(2)

        logger.info("Script 2 executed successfully.")

    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred in Script 2: {str(e)}')

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python RunAll.py <filename>")
        sys.exit(1)

    filename_with_identifier = sys.argv[1]

    process_script1(filename_with_identifier)
    process_script2(filename_with_identifier)

# Save the combined presentation
combined_presentation_path = os.path.join(os.path.splitext(filename_with_identifier)[0] + '_combined.pptx')
combined_presentation.save(combined_presentation_path)
logger.info(f"Combined presentation saved successfully: {combined_presentation_path}")
print(f'Both scripts have been executed and their outputs have been combined into one PowerPoint presentation: "{combined_presentation_path}"')
