import csv
from nltk.tokenize import word_tokenize
from nltk import pos_tag
from pptx import Presentation
from pptx.util import Inches
import os
import sys
import logging
from pptx.enum.text import PP_ALIGN

# Configure logging
logging.basicConfig(filename='separate_verbs.log', level=logging.INFO)

filename_with_identifier = sys.argv[1]
filename_without_extension = os.path.splitext(filename_with_identifier)[0]
file_path = os.path.join(filename_with_identifier)

temp_csv_filename = filename_without_extension + '_temp.csv'

try:
    with open(file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        rows = list(csv_reader)
        # Remove first two rows
        rows = rows[2:]
        # Read the CSV file and extract verbs and their respective cells

        # Extract verbs and non-empty cells from column D
        data = []
        for row in rows:
            for i, cell in enumerate(row):
                if i == 0 or i == 5:
                    continue  

                # Split the cell content into individual words
                words = word_tokenize(cell)
                # Part-of-speech tagging
                tagged_words = pos_tag(words)
                # Check if the cell contains a verb
                if any(tag.startswith('VB') for _, tag in tagged_words):
                    data.append(('Verb', cell))
                # Extract non-empty data from column D
                if i == 3 and cell.strip():  # Check if cell is not empty
                    data.append(('Column D', cell))

        # Write data to temporary CSV file
        with open(temp_csv_filename, 'w', newline='') as temp_csv_file:
            writer = csv.writer(temp_csv_file)
            writer.writerow(['Type', 'Cell'])  # Header
            writer.writerows(data)

    # Create a PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Calculate the positions to fit all cells within the slide
    max_cells_per_row = 5
    num_cells = len(data)
    num_rows = (num_cells + max_cells_per_row - 1) // max_cells_per_row
    cell_width = Inches(2.5)
    cell_height = Inches(0.5)
    left_margin = Inches(0.5)
    top_margin = Inches(1)

    # Create shapes in PowerPoint for data
    left = left_margin
    top = top_margin
    for row in data:
        cell_type, cell_content = row
        if top + cell_height > Inches(10):
            top = top_margin
            left += cell_width
        textbox = slide.shapes.add_textbox(left, top, cell_width, cell_height)
        textbox.text_frame.text = cell_content
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
        top += cell_height

    # Save the PowerPoint presentation
    output_pptx_file = filename_without_extension + '_verbs.pptx'
    prs.save(output_pptx_file)
    print(f"PPTX file with cells containing verbs and non-empty column D data has been created: '{output_pptx_file}'")

except Exception as e:
    # Log any exceptions that occur
    logging.error(f'An error occurred: {str(e)}')
