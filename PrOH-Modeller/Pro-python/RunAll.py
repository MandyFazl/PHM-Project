import csv
import logging
import os
import sys
from nltk.tokenize import word_tokenize
from nltk import pos_tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import spacy
import re

# Create a logger instance
logger = logging.getLogger('merged_logger')
logger.setLevel(logging.INFO)

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('merged_script.log')
file_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)
logger.addHandler(file_handler)

# Load the spaCy English model
nlp = spacy.load("en_core_web_sm")

# Create a single Presentation object
combined_presentation = Presentation()

def process_script1(filename_with_identifier, combined_presentation): #Script1:Core Process Statement
    try:
        logger.info("Processing Script 1")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        logger.info(f"Filename without extension: {filename_without_extension}")

        # Construct the complete file path
        file_path = os.path.join(filename_with_identifier)
        logger.info(f"Uploaded filepath: {file_path}")

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)

        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Define position and size for oval shapes
        num_cells = 7
        oval_width = Inches(2)
        oval_height = Inches(0.8)
        text_font_size = Pt(18)
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

        # Calculate the diagonal line positions
        x_step = slide_width / num_cells
        y_step = slide_height / num_cells

            # Create oval shapes with red background for the first cells in columns A and F
        for i, cell in enumerate(rows[1]):
            # Skip column D
            if i == 3:
                continue
            left = i * x_step
            top = i * y_step
        
            if i == 0:  #First cell in column A
            # Extract the first word from cells A1 and F1
                red_text = rows[1][0].split()[0]  # Get the first word from cell A1
            # Create rectangle shape
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for rectangles

            # Create oval shape with red background
                oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left+x_step, top+y_step, oval_width, oval_height)
                oval_shape.fill.solid()
                oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
    
            elif i == 5:  # First cell in column F
                red_text = rows[1][5].split()[0]  # Get the first word from cell F1
            
            # Create rectangle shape
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left+x_step, top+y_step, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for rectangles
            
            # Create oval shape with red background
                oval_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
                oval_shape.fill.solid()
                oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
            
            elif i == 1: # Column B
            # Create oval shape with green background
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left+x_step, top+y_step, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background
            elif i == 2:  # Column C
            # Create oval shape with red background
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left+x_step, top+y_step, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
            elif i == 4:  # Column E
            # Create oval shape with green background
                shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background

            text_frame = shape.text_frame
            text_frame.text = cell
            text_frame.paragraphs[0].font.size = text_font_size
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Add text to oval shape
            oval_text_frame = oval_shape.text_frame
            oval_text_frame.text = red_text  # Use the extracted first word as text
            oval_text_frame.paragraphs[0].font.size = text_font_size
            oval_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            oval_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Load the spaCy English model
        nlp = spacy.load("en_core_web_sm")
        # Use to extract verbs from the second row, ignoring column D
        verbs = []
        for i, cell in enumerate(rows[1]):
            # Skip column D
            if i == 3:
                continue
            # Apply spaCy NLP pipeline to the cell content
            doc = nlp(cell)
            # Extract verbs
            for token in doc:
                if token.pos_ == 'VERB':
                    verbs.append(token.text)
        # Divide the list of verbs into chunks of five
        verbs_chunks = [verbs[i:i+5] for i in range(0, len(verbs), 5)]
        # Define position and size for text boxes
        left_margin = Inches(0.5)
        top_margin = combined_presentation.slide_height - Inches(0.5)  # Place at the bottom of the slide
        box_width = Inches(1.0)  # Width of each text box
        box_height = Inches(0.5)  # Height of each text box
        # Create text boxes for each chunk of verbs and add them to the bottom of the slide
        for chunk in verbs_chunks:
            left = left_margin
            for word in chunk:
                textbox = slide.shapes.add_textbox(left, top_margin, box_width, box_height)
                textbox.text_frame.text = word
                textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Align the text to the center
                left += box_width
                top_margin -= box_height  # Move to the next line
        # Save the PowerPoint presentation with the same identifier
        # pptx_filename = os.path.join(filename_without_extension + '.pptx')
        # combined_presentation.save(pptx_filename)
        # logger.info("Output presentation is saved successfully.")
        # logger.info(f"Output presentation path: {os.path.join(filename_without_extension + '.pptx')}")
        # print(f'Second row of CSV file has been converted to an editable PowerPoint presentation: "{pptx_filename}"')
    except Exception as e:
        logging.error(f'An error occurred in Script 1: {str(e)}')
        
def process_script2(filename_with_identifier, combined_presentation):  # Script2: Non-Core-Process-Statment
    try:
        # Get the file name with the identifier from the command-line arguments
        filename_with_identifier = sys.argv[1]
        logger.info("looking for the filepath")

        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        logger.info(f"Filename without extension: {filename_without_extension}")

        # Construct the complete file path
        file_path = os.path.join(filename_with_identifier)
        logger.info(f"Uploaded filepath: {file_path}")

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)
            # Remove first two rows
            rows = rows[2:]

        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]  # Blank slide layout
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Define position and size for oval shapes
        num_cells = len(rows[0])
        node_width = Inches(2)
        node_height = Inches(0.8)
        text_font_size = Pt(18)
        top_margin = Inches(1.0)
        left_margin = Inches(0.5)
        left = left_margin
        top = top_margin
        slide_width = combined_presentation.slide_width
        slide_height = combined_presentation.slide_height

    # Calculate the diagonal line positions
        # diagonal_length = (slide_width ** 2 + slide_height ** 2) ** 0.5
        x_step = slide_width / num_cells
        y_step = slide_height / num_cells
    
        # Create oval shapes and add text to them with specified font and background colors
        j = 0
        while j< len(rows):
            row=rows[j]  
            for i, cell in enumerate(row):
                if not cell or i==0 or i==3 or i==5  or '(' in cell: 
                    continue
                if i==1 or i==4:
                    # Create oval shape with Green background
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(0, 255, 0) 
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height
                    
                if i==2:
                    # Create oval shape with Green background
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
                    text_frame = oval.text_frame
                    text_frame.text = cell
                    text_frame.paragraphs[0].font.size = text_font_size
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    left += node_width
                    if left + node_width > Inches(10):
                        left = left_margin
                        top += node_height
            j += 1      
    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred: {str(e)}')
        
def process_script3(filename_with_identifier, combined_presentation): #Script3: SUb-bubbles/ Bracket.py
    try:
        logger.info("Processing Script 3")
        # Input CSV file name
        # filename_with_identifier = sys.argv[1]
        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        # file_path = os.path.join(filename_with_identifier)

        # List to store words between parentheses
        parentheses_words = []
        # Regular expression pattern to find words between parentheses
        pattern = r'\(([^)]*)\)'  # Match everything between '(' and ')'

        with open(filename_with_identifier, 'r', newline='') as csv_input:
            reader = csv.reader(csv_input)
            for row in reader:
                for idx, cell in enumerate(row):
                    # Skip processing cells in columns A and F (indexes 0 and 5 respectively)
                    if idx in [0, 5]:
                        continue
                    # Find all words between parentheses in the cell and store them
                    matches = re.findall(pattern, cell)
                    parentheses_words.extend(matches)
                    
        # Write words between parentheses to a temporary CSV file
        temp_csv_filename = filename_without_extension + '_subbubbles' + '.csv'
        with open(temp_csv_filename, 'w', newline='') as temp_csv_file:
            writer = csv.writer(temp_csv_file)
            writer.writerow(['Words Between Parentheses'])  # Header
            writer.writerows([[word] for word in parentheses_words])   

        # Set to store unique words between parentheses
        unique_parentheses_words = set()

        # Read unique words between parentheses from the temporary CSV file
        with open(temp_csv_filename, 'r', newline='') as csv_input:
            reader = csv.reader(csv_input)
            next(reader)  # Skip the header row

            for row in reader:
                for cell in row:
                    unique_parentheses_words.add(cell)
   
        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]  # Blank slide layout
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Calculate the positions to fit all nodes within the slide
        node_width = Inches(2.5)
        node_height = Inches(1.0)
        left_margin = Inches(0.5)
        top_margin = Inches(1.0)
        left = left_margin
        top = top_margin

        # Create elliptical shapes in PowerPoint for unique words
        for word in unique_parentheses_words:
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            text = shape.text_frame.add_paragraph()
            text.text = word
            text.alignment = 1  # Center align text
            text.font.size = Pt(18)
            text.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
            shape.line.width = Pt(1.5)  # Border width
            left += node_width

            if left + node_width > Inches(10):
                left = left_margin
                top += node_height

        # Save the PowerPoint file
        # pptx_filename = os.path.join(filename_without_extension + '_subbubble' + '.pptx')
        # combined_presentation.save(pptx_filename)
        # logger.info("subbubble is saved successfully.")

    except Exception as e:
        logger.error(f'An error occurred: {str(e)}')
        
def process_script4(filename_with_identifier, combined_presentation): # Script4: Decision_Bubbles.py
    try:
        logger.info("Processing script4")
        filename_without_extension = os.path.splitext(filename_with_identifier)[0]
        logger.info(f"Filename without extension: {filename_without_extension}")

        # Construct the complete file path
        file_path = os.path.join(filename_with_identifier)
        logger.info(f"Uploaded filepath: {file_path}")

        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)

        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]  # Blank slide layout
        slide = combined_presentation.slides.add_slide(slide_layout)

        # Define initial positions for ovals
        left = Inches(0.5)
        top = Inches(1.5)

        # Iterate over all cells in the CSV file
        for row in rows:
            for idx, cell in enumerate(row):
                # Skip processing cells in columns A and F (indexes 0 and 5 respectively)
                if idx in [0, 5]:
                    continue

                # Check if the cell begins with "or:" case-insensitive
                if cell.lower().startswith("or:"):
                    # Extract the phrase after "or:"
                    phrase = cell[3:].strip()

                    # Check if the phrase contains parentheses
                    if '(' in phrase and ')' in phrase:
                        # Split the phrase into main and inner phrases
                        main_phrase, inner_phrase = phrase.split('(', 1)
                        inner_phrase = inner_phrase[:-1]  # Remove the closing parenthesis

                        # Create the outer oval for the main phrase
                        outer_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(3), Inches(1.5))
                        outer_oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                        outer_oval.line.width = Pt(1)  # Border width
                        outer_oval.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dotted line border
                        outer_oval.fill.solid()
                        outer_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                        # Add the main phrase to the outer oval
                        text_frame_outer = outer_oval.text_frame
                        p_outer = text_frame_outer.add_paragraph()
                        p_outer.text = main_phrase.strip()
                        p_outer.font.size = Pt(14)  # Adjust font size as needed
                        p_outer.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                        p_outer.alignment = PP_ALIGN.CENTER  # Center align text

                        # Create the inner oval for the phrase inside parentheses
                        inner_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.5), top + Inches(0.3), Inches(2), Inches(0.8))
                        inner_oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                        inner_oval.line.width = Pt(1)  # Border width
                        inner_oval.fill.solid()
                        inner_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                        # Add the inner phrase to the inner oval
                        text_frame_inner = inner_oval.text_frame
                        p_inner = text_frame_inner.add_paragraph()
                        p_inner.text = inner_phrase.strip()
                        p_inner.font.size = Pt(12)  # Adjust font size as needed
                        p_inner.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                        p_inner.alignment = PP_ALIGN.CENTER  # Center align text

                    else:
                        # Create a black dotted line oval for the extracted phrase
                        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(3), Inches(1.5))
                        oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                        oval.line.width = Pt(1)  # Border width
                        oval.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dotted line border
                        oval.fill.solid()
                        oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                        # Add the extracted phrase to the oval
                        text_frame = oval.text_frame
                        p = text_frame.add_paragraph()
                        p.text = phrase.strip()
                        p.font.size = Pt(14)  # Adjust font size as needed
                        p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                        p.alignment = PP_ALIGN.CENTER  # Center align text

                    # Update positions for the next oval
                    left += Inches(3.5)
                    if left + Inches(3) > Inches(10):
                        left = Inches(0.5)
                        top += Inches(2)

        logger.info("script3 executed successfully.")

    except Exception as e:
        # Log any exceptions that occur
        logging.error(f'An error occurred in Script 2: {str(e)}')

def process_script5(filename_with_identifier, combined_presentation): # Script5: Seperate_verbs.py
    filename_without_extension = os.path.splitext(filename_with_identifier)[0]
    file_path = os.path.join(filename_with_identifier)
    temp_csv_filename = filename_without_extension + '_verbs.csv'
    try:
        with open(file_path, 'r') as csv_file:
            csv_reader = csv.reader(csv_file)
            rows = list(csv_reader)

            # Remove first two rows
            rows = rows[2:]
            # Read the CSV file and extract verbs and their respective cells
            # Extract verbs and non-empty cells from column D
            data = []
            for row in rows:
                for i, cell in enumerate(row):
                    if i == 0 or i == 5:
                        continue   
                    # Split the cell content into individual words
                    words = word_tokenize(cell)
                    # Part-of-speech tagging
                    tagged_words = pos_tag(words)
                    # Check if the cell contains a verb
                    if any(tag.startswith('VB') for _, tag in tagged_words):
                        data.append(('Verb', cell))
                    # Extract non-empty data from column D
                    if i == 3 and cell.strip():  # Check if cell is not empty
                        data.append(('Column D', cell))
        # Create a blank slide
        slide_layout = combined_presentation.slide_layouts[5]  # Blank slide layout
        slide = combined_presentation.slides.add_slide(slide_layout)
            
        # Calculate the positions to fit all cells within the slide
        max_cells_per_row = 5
        num_cells = len(data)
        num_rows = (num_cells + max_cells_per_row - 1) // max_cells_per_row
        cell_width = Inches(2.5)
        cell_height = Inches(0.5)
        left_margin = Inches(0.5)
        top_margin = Inches(1)

        # Create shapes in PowerPoint for data
        left = left_margin
        top = top_margin
        for row in data:
            cell_type, cell_content = row
            if top + cell_height > Inches(10):
                top = top_margin
                left += cell_width
            textbox = slide.shapes.add_textbox(left, top, cell_width, cell_height)
            textbox.text_frame.text = cell_content
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
            top += cell_height
        
    except Exception as e:
        logging.error(f'An error occurred: {str(e)}')


    # Main function
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python RunAll.py <filename>")
        sys.exit(1)

    filename_with_identifier = sys.argv[1]

    process_script1(filename_with_identifier, combined_presentation)
    process_script2(filename_with_identifier, combined_presentation)
    process_script3(filename_with_identifier, combined_presentation)
    process_script4(filename_with_identifier, combined_presentation)
    process_script5(filename_with_identifier, combined_presentation)

    # Save the combined presentation
    combined_presentation_path = os.path.join(os.path.splitext(filename_with_identifier)[0] + '_combined.pptx')
    combined_presentation.save(combined_presentation_path)
    logger.info(f"Combined presentation saved successfully: {combined_presentation_path}")
    print(f'All scripts have been executed and their outputs have been combined into one PowerPoint presentation: "{combined_presentation_path}"')

# the first 3 scripts works well every time you want to execute this code run the following:python RunAll.py uploads\make_car_d27a592c.csv
#after other alternation I will get back and finalise this code                  
