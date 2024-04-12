import csv
import os
import sys
import logging
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new logger instance for bracket.py
logger = logging.getLogger('bracket_logger')
logger.setLevel(logging.INFO)

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('bracket.log')
file_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)
logger.addHandler(file_handler)

try:
    # Input CSV file name
    filename_with_identifier = sys.argv[1]
    filename_without_extension = os.path.splitext(filename_with_identifier)[0]
    file_path = os.path.join(filename_with_identifier)

    # List to store words between parentheses
    parentheses_words = []

    # Regular expression pattern to find words between parentheses
    pattern = r'\(([^)]*)\)'  # Match everything between '(' and ')'

    with open(file_path, 'r', newline='') as csv_input:
        reader = csv.reader(csv_input)

        for row in reader:
            for idx, cell in enumerate(row):
                # Skip processing cells in columns A and F (indexes 0 and 5 respectively)
                if idx in [0, 5]:
                    continue
                
                # Find all words between parentheses in the cell and store them
                matches = re.findall(pattern, cell)
                parentheses_words.extend(matches)

    # Write words between parentheses to a temporary CSV file
    temp_csv_filename = filename_without_extension + '_subbubbles' + '.csv'
    with open(temp_csv_filename, 'w', newline='') as temp_csv_file:
        writer = csv.writer(temp_csv_file)
        writer.writerow(['Words Between Parentheses'])  # Header
        writer.writerows([[word] for word in parentheses_words])   
    logger.info(f"Words between parentheses have been saved to '{temp_csv_filename}'")

    # Set to store unique words between parentheses
    unique_parentheses_words = set()

    # Read unique words between parentheses from the temporary CSV file
    with open(temp_csv_filename, 'r', newline='') as csv_input:
        reader = csv.reader(csv_input)
        next(reader)  # Skip the header row

        for row in reader:
            for cell in row:
                unique_parentheses_words.add(cell)

    # Create a PowerPoint presentation
    Presentation = Presentation()
    slide_layout = Presentation.slide_layouts[5]  # Use a blank slide layout

    # Create a single slide with a white background
    slide = Presentation.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

    # Calculate the positions to fit all nodes within the slide
    num_nodes = len(unique_parentheses_words)
    node_width = Inches(2.5)
    node_height = Inches(1.0)
    left_margin = Inches(0.5)
    top_margin = Inches(1.0)
    left = left_margin
    top = top_margin

    # Create elliptical shapes in PowerPoint for unique words
    for word in unique_parentheses_words:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, node_width, node_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        text = shape.text_frame.add_paragraph()
        text.text = word
        text.alignment = 1  # Center align text
        text.font.size = Pt(18)
        text.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
        shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        shape.line.width = Pt(1.5)  # Border width
        left += node_width

        if left + node_width > Inches(10):
            left = left_margin
            top += node_height

    # Save the PowerPoint file
    output_pptx_file =os.path.join(filename_without_extension +'_subbubbles'+'.pptx')
    Presentation.save(output_pptx_file)
    logger.info(f"PPTX file with editable elliptical nodes for unique values has been created: '{output_pptx_file}'")

except Exception as e:
    # Log any exceptions that occur
    logger.error(f'An error occurred: {str(e)}')
